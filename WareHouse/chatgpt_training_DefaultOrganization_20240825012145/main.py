'''
Main file for the ChatGPT PowerPoint Creator
'''
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
def create_powerpoint():
    # Create a new PowerPoint presentation
    presentation = Presentation()
    # Get the input text from the user
    input_text = text_entry.get("1.0", tk.END)
    # Split the input text into slides
    slides = input_text.split("\n---\n")
    # Create slides in the PowerPoint presentation
    for slide_text in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = slide_text.split("\n")[0]
        content.text = "\n".join(slide_text.split("\n")[1:])
    # Save the PowerPoint presentation
    save_path = filedialog.asksaveasfilename(defaultextension=".pptx")
    if save_path:
        presentation.save(save_path)
        messagebox.showinfo("Success", "PowerPoint presentation created successfully!")
# Create the main window
window = tk.Tk()
window.title("ChatGPT PowerPoint Creator")
# Create a text entry field for the user to input the slides
text_entry = tk.Text(window, height=10, width=50)
text_entry.pack()
# Create a button to create the PowerPoint presentation
create_button = tk.Button(window, text="Create PowerPoint", command=create_powerpoint)
create_button.pack()
# Start the main event loop
window.mainloop()